from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Bengaluru 3D Digital Twin"
    slide.placeholders[1].text = "Nexus Twin: Admin & Citizen Interface\nThe future of urban governance is in pixels."

    # Slide 2: Admin Nexus Features (Part 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Admin Nexus: Strategic Command (1/2)"
    tf = slide.placeholders[1].text_frame
    tf.text = "3D Digital Twin Visualization: High-fidelity rendering of Bengaluru's infrastructure."
    p = tf.add_paragraph()
    p.text = "Agent-Based Simulation (ABM): Real-time simulation of 400+ autonomous agents."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Deck.gl TripsLayer: Neon motion trails for real-time traffic flow."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Impact Audit Simulation: 'Demolish Mode' to analyze structural changes."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Predictive Failure Projections: Storm intensity-based risk mapping."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "AI Policy Advisor: Natural language interface for urban governance."
    p.level = 0

    # Slide 3: Admin Nexus Features (Part 2)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Admin Nexus: Strategic Command (2/2)"
    tf = slide.placeholders[1].text_frame
    tf.text = "'The Sims' Asset Deployment: Drag-and-drop system for placing 3D assets."
    p = tf.add_paragraph()
    p.text = "Triple Win Scorecard: Real-time KPI tracking for Economic, Social, Environmental performance."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Crisis Command Center: Interactive flood inundation modeling."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Social Sentiment Analytics: Aggregated citizen mood heatmaps."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "X-Ray Utility Vision: Subsurface visualization of networks."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Strategic Activity Log: Real-time terminal feed of administrative events."
    p.level = 0

    # Slide 4: Citizen Nexus Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Citizen Nexus: Public Observatory"
    tf = slide.placeholders[1].text_frame
    tf.text = "Public Urban Observatory: Read-only 3D view of city infrastructure."
    p = tf.add_paragraph()
    p.text = "Sentiment Heatmap Access: Visualization of city-wide public mood."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Crowdsourced Issue Reporting: Map-based targeting to pin local issues."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Environmental Health HUD: Toggles for AQI and vegetation density layers."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Public Safety Flood Simulator: Citizen-facing disaster awareness models."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Heritage Timeline Engine: Historical architectural slider (1920 to 2024)."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Live Traffic Pulse & Atmospheric Weather Engine."
    p.level = 0

    prs.save('Bengaluru_Nexus_Twin_Presentation_V2.pptx')

if __name__ == '__main__':
    create_presentation()
